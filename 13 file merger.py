"""
Generates a merged word, ppt, pdf file of 2 or all files given in the folder path.

"""
import os
import PyPDF2
from pptx import Presentation
from docx import Document

def merge_two_pdfs(file1_path, file2_path, output_file):
    """Merges two PDF files into a single output file."""
    try:
        pdf_writer = PyPDF2.PdfWriter()
        
        for file_path in [file1_path, file2_path]:
            with open(file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)
        
        with open(output_file, 'wb') as output_pdf:
            pdf_writer.write(output_pdf)
        
        print("PDFs merged successfully!")
    except Exception as e:
        print(f"Error merging PDFs: {e}")

def merge_all_pdfs(folder_path, output_file):
    """Merges all PDF files in the given folder into a single PDF."""
    try:
        pdf_files = [f for f in os.listdir(folder_path) if f.endswith('.pdf')]
        
        if not pdf_files:
            print("No PDF files found in the folder.")
            return
        
        print("Merging the following PDFs:")
        for pdf in pdf_files:
            print(pdf)
        
        pdf_writer = PyPDF2.PdfWriter()
        
        for pdf_file in pdf_files:
            pdf_path = os.path.join(folder_path, pdf_file)
            with open(pdf_path, 'rb') as pdf:
                pdf_reader = PyPDF2.PdfReader(pdf)
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)
        
        with open(output_file, 'wb') as output_pdf:
            pdf_writer.write(output_pdf)
        
        print(f"All PDFs merged successfully into {output_file}")
    except Exception as e:
        print(f"Error merging PDFs: {e}")

def merge_two_ppts(file1_path, file2_path, output_file):
    """Merges two PPT files into a single presentation."""
    try:
        merged_ppt = Presentation()
        
        for file_path in [file1_path, file2_path]:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                new_slide = merged_ppt.slides.add_slide(merged_ppt.slide_layouts[5])
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        new_text_box = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                        new_text_box.text = shape.text
        
        merged_ppt.save(output_file)
        print("PPTs merged successfully!")
    except Exception as e:
        print(f"Error merging PPTs: {e}")

def merge_all_ppts(folder_path, output_file):
    """Merges all PPT files in the given folder into a single presentation."""
    try:
        ppt_files = [f for f in os.listdir(folder_path) if f.endswith('.pptx')]
        
        if not ppt_files:
            print("No PPT files found in the folder.")
            return
        
        print("Merging the following PPTs:")
        for ppt in ppt_files:
            print(ppt)
        
        merged_ppt = Presentation()
        
        for ppt_file in ppt_files:
            ppt_path = os.path.join(folder_path, ppt_file)
            presentation = Presentation(ppt_path)
            for slide in presentation.slides:
                new_slide = merged_ppt.slides.add_slide(merged_ppt.slide_layouts[5])
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        new_text_box = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                        new_text_box.text = shape.text
        
        merged_ppt.save(output_file)
        print(f"All PPTs merged successfully into {output_file}")
    except Exception as e:
        print(f"Error merging PPTs: {e}")

def merge_two_word_docs(file1_path, file2_path, output_file):
    """Merges two Word documents into a single document."""
    try:
        merged_doc = Document()
        
        for file_path in [file1_path, file2_path]:
            doc = Document(file_path)
            for para in doc.paragraphs:
                merged_doc.add_paragraph(para.text)
        
        merged_doc.save(output_file)
        print("Word documents merged successfully!")
    except Exception as e:
        print(f"Error merging Word documents: {e}")

def merge_all_word_docs(folder_path, output_file):
    """Merges all Word documents in the given folder into a single document."""
    try:
        word_files = [f for f in os.listdir(folder_path) if f.endswith('.docx')]
        
        if not word_files:
            print("No Word files found in the folder.")
            return
        
        print("Merging the following Word files:")
        for doc in word_files:
            print(doc)
        
        merged_doc = Document()
        
        for word_file in word_files:
            doc_path = os.path.join(folder_path, word_file)
            doc = Document(doc_path)
            for para in doc.paragraphs:
                merged_doc.add_paragraph(para.text)
        
        merged_doc.save(output_file)
        print(f"All Word documents merged successfully into {output_file}")
    except Exception as e:
        print(f"Error merging Word documents: {e}")

def main():
    folder_path = input("Enter the folder path containing files: ")
    choice = input("Do you want to merge PDFs, PPTs, or Word files? (pdf/ppt/doc): ").strip().lower()
    
    if choice == "pdf":
        output_file = os.path.join(folder_path, "merged_all.pdf")
        merge_all_pdfs(folder_path, output_file)
    elif choice == "ppt":
        output_file = os.path.join(folder_path, "merged_presentation.pptx")
        merge_all_ppts(folder_path, output_file)
    elif choice == "doc":
        output_file = os.path.join(folder_path, "merged_document.docx")
        merge_all_word_docs(folder_path, output_file)
    else:
        print("Invalid choice. Please enter 'pdf', 'ppt', or 'doc'.")

if __name__ == "__main__":
    main()
